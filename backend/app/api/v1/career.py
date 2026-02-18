"""
Career analytics API endpoints
"""
import os
import re
import uuid
from pathlib import Path
from fastapi import APIRouter, HTTPException, UploadFile, File, Depends, Query
from typing import List, Optional
from pydantic import BaseModel
from sqlalchemy import select, delete
from sqlalchemy.ext.asyncio import AsyncSession
from app.services.career.analytics_service import CareerAnalytics
from app.services.career.company_suggestion_service import CompanySuggestionService

try:
    from app.services.career.rag_service import RAGEngine
    _rag_engine = RAGEngine()
except Exception:
    _rag_engine = None
from app.services.llm.llm_service import LLMService
from app.services.resume.parser import ResumeParser
from app.config import settings
from app.database.base import get_db
from app.database.models.pathfinder import (
    PathfinderUser,
    PathfinderUserProfile,
    PathfinderUserDocument,
    PathfinderUserCoursework,
    PathfinderUserProject,
    PathfinderUserCareerInterest,
    PathfinderJob,
)

router = APIRouter()
resume_parser = ResumeParser()

# Labels from transcript PDFs that are not course names (Student Information block, etc.)
_NON_COURSE_LABELS = frozenset({
    "student information", "student id", "student name", "phone", "address", "advisor",
    "degree", "term", "cumulative", "gpa group", "graduate", "undergraduate",
    "course", "course name", "grade", "credits", "subject", "code",
    "attempted", "earned", "hours", "grade points", "repeat",
    "—", "-", ""
})


def _is_course_row(course: str) -> bool:
    """Return False if this looks like a Student Info label or header, not a real course."""
    if not course or len(course) < 2:
        return False
    key = course.lower().strip()
    # Strip trailing colon and use the part before colon for label check (e.g. "Student ID:" -> "student id")
    if ":" in key:
        label_part = key.split(":")[0].strip()
        if label_part in _NON_COURSE_LABELS:
            return False
        # Reject "GPA Group: Graduate" style rows
        if label_part == "gpa group" or key.startswith("gpa group"):
            return False
    if key in _NON_COURSE_LABELS:
        return False
    if key.startswith(("student id", "phone", "degree", "address", "advisor", "cumulative", "attempted", "earned", "gpa group", "graduate", "undergraduate")):
        return False
    if key.replace(".", "").replace(" ", "").isdigit():
        return False
    return True


def _filter_course_rows(course_grades: list) -> list:
    """Remove rows that are Student Information or other non-course entries."""
    return [r for r in (course_grades or []) if _is_course_row(r.get("course"))]


def _looks_like_grade(val: str) -> bool:
    """True if value looks like a letter grade (A, B+, IP, A-) or numeric grade (e.g. 85)."""
    if not val or not val.strip():
        return False
    import re
    v = val.strip().upper()
    if re.match(r"^[A-F][+-]?$", v, re.IGNORECASE):
        return True
    if v == "IP":  # In Progress
        return True
    if re.match(r"^\d{1,3}(\.\d+)?$", v):
        return True
    return False


def _looks_like_credits(val: str) -> bool:
    """True if value looks like credit hours (e.g. 3, 1, 4.0)."""
    if not val or not val.strip():
        return False
    import re
    v = val.strip()
    if not re.match(r"^\d{1,2}(\.\d+)?$", v):
        return False
    try:
        n = float(v)
        return 0.5 <= n <= 15
    except ValueError:
        return False


def _find_credits_grade_columns(cells: list) -> tuple:
    """
    Given a header row (list of cell strings), return (credits_col, grade_col) 0-based indices, or (None, None).
    Handles 'Credits' before or after 'Grade' and common variants.
    """
    credits_col, grade_col = None, None
    for i, c in enumerate(cells):
        if not c:
            continue
        k = str(c).strip().lower()
        if k in ("credits", "credit", "units", "hrs", "ch"):
            credits_col = i
        if k in ("grade", "grades", "letter", "score"):
            grade_col = i
    return (credits_col, grade_col)


def _assign_grade_credits_from_cells(cells: list, start_idx: int, credits_col: int, grade_col: int) -> tuple:
    """
    From a data row cells, get (grade, credits) using either known column indices or format detection.
    start_idx: first data column index (after course/code); we look at cells from start_idx.
    """
    grade, credits = None, None
    # Use header indices if we have them and they're in range
    if credits_col is not None and credits_col < len(cells):
        v = (cells[credits_col] or "").strip()
        if _looks_like_credits(v):
            credits = v
    if grade_col is not None and grade_col < len(cells):
        v = (cells[grade_col] or "").strip()
        if _looks_like_grade(v) or v == "IP":
            grade = v
    # If we got both from header, done
    if grade is not None and credits is not None:
        return (grade, credits)
    # Otherwise scan cells from start_idx and assign by format (credits = small number, grade = letter/IP)
    for j in range(start_idx, min(len(cells), start_idx + 6)):
        if j >= len(cells):
            break
        v = (cells[j] or "").strip()
        if not v:
            continue
        if _looks_like_credits(v) and credits is None:
            credits = v
        if (_looks_like_grade(v) or v.upper() == "IP") and grade is None:
            grade = v
    return (grade or None, credits or None)


def _extract_course_grades_from_pdf_tables(file_path: str) -> list:
    """
    Extract course, credits, and grade from PDF tables (e.g. SFBU transcript).
    Uses header row when present (Credits / Grade columns); else detects by format.
    Returns list of dicts: [{"course", "grade", "credits"}].
    """
    try:
        import pdfplumber
        out = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                tables = page.extract_tables()
                for table in (tables or []):
                    if not table or len(table) < 2:
                        continue
                    credits_col, grade_col = None, None
                    header_skipped = False
                    for i, row in enumerate(table):
                        if not row or not any(cell and str(cell).strip() for cell in row):
                            continue
                        cells = [str(c or "").strip() for c in row]
                        # Detect header: row containing "Credits" and/or "Grade"
                        if i == 0 or not header_skipped:
                            cr_idx, gr_idx = _find_credits_grade_columns(cells)
                            if cr_idx is not None or gr_idx is not None:
                                credits_col, grade_col = cr_idx, gr_idx
                                header_skipped = True
                                if cells and cells[0].lower() in ("course", "course name", "subject", "code", "credits", "grade"):
                                    continue
                        if i == 0 and cells and cells[0].lower() in ("course", "course name", "subject", "code"):
                            continue
                        # Course: first column is often code (CS501) or combined code+name
                        course = (cells[0] or "").strip() or None
                        if not _is_course_row(course):
                            continue
                        # If we have a second column that looks like a long title and first is short code, use both for course
                        if len(cells) >= 2 and cells[1] and _looks_like_credits(cells[1]) is False and _looks_like_grade(cells[1]) is False:
                            maybe_name = (cells[1] or "").strip()
                            if len(maybe_name) > len(course or "") and " - " in maybe_name:
                                course = maybe_name  # e.g. "DS512 - Data Engineering"
                        start_idx = 1 if course == (cells[0] or "").strip() else 2
                        grade, credits = _assign_grade_credits_from_cells(cells, start_idx, credits_col, grade_col)
                        out.append({"course": course or "", "grade": grade, "credits": credits})
        return _filter_course_rows(out)[:120]
    except Exception:
        return []


# Initialize services
career_analytics = CareerAnalytics()
company_suggestion_service = CompanySuggestionService()
llm_service = LLMService()


class CareerInsightsRequest(BaseModel):
    """Request model for career insights"""
    resume_text: str
    user_query: Optional[str] = None


class SkillDemandRequest(BaseModel):
    """Request model for skill demand analysis"""
    user_skills: List[str]


class SalaryInsightsRequest(BaseModel):
    """Request model for salary insights"""
    user_skills: List[str]
    experience_level: Optional[int] = None


class CompanySuggestionsRequest(BaseModel):
    """Request model for company suggestions from coursework, projects, interests"""
    coursework: Optional[List[str]] = None
    projects: Optional[List[str]] = None
    interests: Optional[List[str]] = None
    target_role: Optional[str] = None
    limit: int = 10


class ExtractCoursesRequest(BaseModel):
    """Request model for extracting course names from pasted text (e.g. SFBU course-grades)"""
    raw_text: str


class ImportCourseGradesRequest(BaseModel):
    """Request model for importing course grades from pasted SFBU content"""
    raw_text: str


class AnalyzeCourseworkRequest(BaseModel):
    """Request model for analyzing coursework, resume, and projects."""
    course_grades: List[dict]  # [{"course": str, "grade": str|None, "credits": str|None}]
    resume_text: Optional[str] = None
    projects: Optional[List[str]] = None
    job_area_interest: Optional[str] = None


@router.post("/career/insights")
async def get_career_insights(request: CareerInsightsRequest):
    """
    Get comprehensive career insights using RAG
    
    Args:
        request: Career insights request
        
    Returns:
        Career insights data
    """
    if _rag_engine is None:
        raise HTTPException(
            status_code=503,
            detail="Career insights (RAG) unavailable on this build. Use full requirements for full features."
        )
    try:
        insights = await _rag_engine.get_career_insights(
            request.resume_text,
            request.user_query
        )
        return insights
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/career/skill-demand")
async def get_skill_demand(request: SkillDemandRequest):
    """
    Analyze skill demand in job market
    
    Args:
        request: Skill demand request
        
    Returns:
        Skill demand analysis
    """
    try:
        analysis = career_analytics.get_skill_demand_analysis(request.user_skills)
        return analysis
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/career/salary-insights")
async def get_salary_insights(request: SalaryInsightsRequest):
    """
    Get salary insights based on skills and experience
    
    Args:
        request: Salary insights request
        
    Returns:
        Salary statistics
    """
    try:
        insights = career_analytics.get_salary_insights(
            request.user_skills,
            request.experience_level
        )
        return insights
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/career/industry-insights")
async def get_industry_insights(request: SkillDemandRequest):
    """
    Get industry insights based on skills
    
    Args:
        request: Skill demand request
        
    Returns:
        Industry distribution
    """
    try:
        insights = career_analytics.get_industry_insights(request.user_skills)
        return insights
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/career/extract-courses")
async def extract_courses(request: ExtractCoursesRequest):
    """
    Extract course names from pasted text (e.g. from SFBU course-grades page).
    Returns a list of course names for use in company suggestions.
    """
    try:
        courses = await llm_service.extract_courses_from_text(request.raw_text)
        return {"courses": courses}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/career/import-course-grades")
async def import_course_grades(request: ImportCourseGradesRequest):
    """
    Import and parse course grades from pasted SFBU course-grades page content.
    Returns structured list: [{ "course", "grade", "credits" }] for display and analysis.
    """
    try:
        course_grades = await llm_service.extract_course_grades_from_text(request.raw_text)
        if not course_grades:
            course_names = await llm_service.extract_courses_from_text(request.raw_text)
            course_grades = [{"course": c, "grade": None, "credits": None} for c in course_names]
        else:
            course_grades = llm_service.fill_grades_credits_from_text(course_grades, request.raw_text)
        return {"course_grades": course_grades or []}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/career/import-course-grades-pdf")
async def import_course_grades_pdf(file: UploadFile = File(...)):
    """
    Import course grades from an exported SFBU transcript/course-grades PDF.
    Export the PDF from the SFBU portal, upload here; we extract text and parse courses and grades.
    """
    if not file.filename or not file.filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Please upload a PDF file (e.g. exported from SFBU course grades).")
    file_id = str(uuid.uuid4())
    save_path = settings.UPLOAD_DIR / f"transcript_{file_id}.pdf"
    try:
        content = await file.read()
        if len(content) > settings.MAX_UPLOAD_SIZE:
            raise HTTPException(status_code=400, detail="File too large.")
        with open(save_path, "wb") as f:
            f.write(content)
        # Extract text: try pdfplumber for tables first (better for grades), else pypdf
        raw_text = await resume_parser.parse_resume(str(save_path), preserve_case=True)
        if not raw_text or len(raw_text.strip()) < 20:
            raise HTTPException(
                status_code=400,
                detail="Could not extract enough text from the PDF. Try exporting again or use a PDF with selectable text.",
            )
        course_grades = _extract_course_grades_from_pdf_tables(str(save_path))
        if not course_grades:
            course_grades = await llm_service.extract_course_grades_from_text(raw_text)
            course_grades = llm_service.fill_grades_credits_from_text(course_grades, raw_text)
        if not course_grades:
            course_names = await llm_service.extract_courses_from_text(raw_text)
            course_grades = [{"course": c, "grade": None, "credits": None} for c in course_names]
        course_grades = _filter_course_rows(course_grades or [])
        return {
            "course_grades": course_grades,
            "extracted_text_preview": raw_text[:500],
            "extracted_text": raw_text[:8000],
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        if save_path.exists():
            try:
                os.remove(save_path)
            except Exception:
                pass


def _extract_text_from_pdf_plumber(file_path: str) -> str:
    """Fallback PDF text extraction using pdfplumber (often better for complex layouts)."""
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t and t.strip():
                    parts.append(t.strip())
        return "\n".join(parts).strip()
    except Exception:
        return ""


def _extract_shape_text(shape) -> str:
    """Recursively extract text from a shape (handles groups, tables)."""
    parts = []
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                t = _extract_shape_text(child)
                if t:
                    parts.append(t)
        else:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text and para.text.strip():
                        parts.append(para.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text and cell.text.strip():
                            parts.append(cell.text.strip())
    except Exception:
        pass
    return "\n".join(p for p in parts if p)


def _extract_text_from_pptx_xml(file_path: str) -> str:
    """Fallback: extract text from PPTX as ZIP + XML when python-pptx returns empty."""
    import zipfile
    import xml.etree.ElementTree as ET
    parts = []
    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            for name in sorted(zf.namelist()):
                if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                    try:
                        with zf.open(name) as f:
                            root = ET.parse(f).getroot()
                            for elem in root.iter():
                                if elem.tag.endswith("}t"):
                                    if elem.text:
                                        t = elem.text.strip()
                                        if t:
                                            parts.append(t)
                                    for child in elem:
                                        if child.tail and child.tail.strip():
                                            parts.append(child.tail.strip())
                    except Exception:
                        pass
                if name.startswith("ppt/notesSlides/notesSlide") and name.endswith(".xml"):
                    try:
                        with zf.open(name) as f:
                            root = ET.parse(f).getroot()
                            for elem in root.iter():
                                if elem.tag.endswith("}t"):
                                    if elem.text:
                                        t = elem.text.strip()
                                        if t:
                                            parts.append(t)
                    except Exception:
                        pass
        return "\n".join(parts).strip() if parts else ""
    except Exception:
        return ""


def _extract_text_from_pptx(file_path: str) -> str:
    """Extract text from all slides, notes, and shapes in a .pptx file."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                t = _extract_shape_text(shape)
                if t:
                    parts.append(t)
            # Extract slide notes
            try:
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes and notes.strip():
                        parts.append(notes.strip())
            except Exception:
                pass
        text = "\n".join(p for p in parts if p).strip()
        if text:
            return text
        # Fallback: extract from PPTX as ZIP + XML (handles shapes python-pptx misses)
        return _extract_text_from_pptx_xml(file_path)
    except Exception:
        return _extract_text_from_pptx_xml(file_path)


def _extract_text_from_docx(file_path: str) -> str:
    """Extract text from a .docx file. Tries docx2txt first (pure Python), then python-docx."""
    # Try docx2txt first - lightweight, no external deps
    try:
        import docx2txt
        text = docx2txt.process(file_path)
        if text and text.strip():
            return text.strip()
    except Exception:
        pass
    # Fallback to python-docx
    try:
        from docx import Document
        doc = Document(file_path)
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        parts.append(cell.text)
        return "\n".join(parts).strip()
    except Exception:
        return ""


@router.post("/career/import-project-files")
async def import_project_files(files: List[UploadFile] = File(...)):
    """
    Import multiple project files (PDF, PPTX, DOCX). Extracts text from each and returns
    list of {filename, text} for use in analyse. Max 20 files.
    """
    if len(files) > 20:
        raise HTTPException(status_code=400, detail="Maximum 20 project files.")
    allowed = (".pdf", ".pptx", ".docx")
    results = []
    for u in files:
        if not u.filename:
            continue
        low = u.filename.lower()
        if not any(low.endswith(ext) for ext in allowed):
            results.append({"filename": u.filename, "text": "", "error": f"Unsupported format. Use {', '.join(allowed)}"})
            continue
        file_id = str(uuid.uuid4())
        if low.endswith(".pdf"):
            ext = ".pdf"
        elif low.endswith(".docx"):
            ext = ".docx"
        else:
            ext = ".pptx"
        save_path = settings.UPLOAD_DIR / f"project_{file_id}{ext}"
        try:
            content = await u.read()
            if len(content) > settings.MAX_UPLOAD_SIZE:
                results.append({"filename": u.filename, "text": "", "error": "File too large."})
                continue
            with open(save_path, "wb") as f:
                f.write(content)
            if ext == ".pdf":
                text = await resume_parser.parse_resume(str(save_path), preserve_case=True)
                if not (text or "").strip():
                    text = _extract_text_from_pdf_plumber(str(save_path))
            elif ext == ".docx":
                text = _extract_text_from_docx(str(save_path))
            else:
                text = _extract_text_from_pptx(str(save_path))
            text = (text or "").strip()
            results.append({"filename": u.filename, "text": text[:15000], "error": None})
        except HTTPException:
            raise
        except Exception as e:
            results.append({"filename": u.filename, "text": "", "error": str(e)})
        finally:
            if save_path.exists():
                try:
                    os.remove(save_path)
                except Exception:
                    pass
    return {"projects": results}


@router.post("/career/extract-resume-pdf")
async def extract_resume_pdf(file: UploadFile = File(...)):
    """
    Extract text from a resume PDF. Use this before analyze-coursework to analyze based on coursework + resume.
    """
    if not file.filename or not file.filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Please upload a PDF file (resume).")
    file_id = str(uuid.uuid4())
    save_path = settings.UPLOAD_DIR / f"resume_{file_id}.pdf"
    try:
        content = await file.read()
        if len(content) > settings.MAX_UPLOAD_SIZE:
            raise HTTPException(status_code=400, detail="File too large.")
        with open(save_path, "wb") as f:
            f.write(content)
        raw_text = await resume_parser.parse_resume(str(save_path), preserve_case=True)
        if not raw_text or len(raw_text.strip()) < 10:
            raise HTTPException(status_code=400, detail="Could not extract enough text from the PDF.")
        return {"resume_text": raw_text}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        if save_path.exists():
            try:
                os.remove(save_path)
            except Exception:
                pass


@router.post("/career/analyze-coursework")
async def analyze_coursework(request: AnalyzeCourseworkRequest):
    """
    Analyze imported coursework (and grades). Optionally include resume_text to analyze based on both.
    Returns summary, strengths, suggested roles, skills to highlight, and recommendations.
    """
    try:
        result = await llm_service.analyze_coursework(
            request.course_grades,
            resume_text=request.resume_text,
            projects=request.projects,
            job_area_interest=request.job_area_interest,
        )
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


class ExtractProfileRequest(BaseModel):
    """Request model for extracting profile from documents."""
    resume_text: Optional[str] = None
    course_grades: Optional[List[dict]] = None
    coursework_raw_text: Optional[str] = None
    projects: Optional[List[str]] = None


@router.post("/career/extract-profile")
async def extract_profile(request: ExtractProfileRequest):
    """
    Extract profile (name, academic title, skills, courses with term and tags) from resume, coursework, and projects.
    Use after uploading documents to populate the dashboard.
    """
    try:
        result = await llm_service.extract_profile(
            resume_text=request.resume_text,
            course_grades=request.course_grades or [],
            coursework_raw_text=request.coursework_raw_text,
            projects=request.projects,
        )
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/career/company-suggestions")
async def get_company_suggestions(request: CompanySuggestionsRequest):
    """
    Analyze student's coursework, projects, and interests
    to suggest specific companies that match their profile.
    """
    try:
        result = await company_suggestion_service.suggest_companies(
            coursework=request.coursework,
            projects=request.projects,
            interests=request.interests,
            target_role=request.target_role,
            limit=request.limit,
        )
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/career/roadmap")
async def generate_roadmap(request: CareerInsightsRequest):
    """
    Generate career roadmap
    
    Args:
        request: Career insights request
        
    Returns:
        Career roadmap
    """
    try:
        roadmap = await llm_service.generate_career_roadmap(
            request.resume_text
        )
        return {"roadmap": roadmap}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/career/jobs/search")
async def search_jobs(query: str, limit: int = 10):
    """
    Search for relevant jobs
    
    Args:
        query: Search query
        limit: Number of results
        
    Returns:
        List of relevant jobs
    """
    if _rag_engine is None:
        raise HTTPException(
            status_code=503,
            detail="Job search (RAG) unavailable on this build. Use full requirements for full features."
        )
    try:
        jobs = _rag_engine.search_relevant_jobs(query, n_results=limit)
        return {"jobs": jobs}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# --- Profile persistence & related jobs ---

class SaveProfileRequest(BaseModel):
    """Save extracted profile to PostgreSQL."""
    email: str
    name: Optional[str] = None
    academic_title: Optional[str] = None
    technical_skills: Optional[List[dict]] = None
    soft_skills: Optional[List[dict]] = None
    courses: Optional[List[dict]] = None
    profile_projects: Optional[List[dict]] = None
    documents: Optional[List[dict]] = None
    career_interests: Optional[List[str]] = None


@router.post("/career/save-profile")
async def save_profile(request: SaveProfileRequest, db: AsyncSession = Depends(get_db)):
    """Save or update user profile in PostgreSQL."""
    email = (request.email or "").strip().lower()
    if not email:
        raise HTTPException(status_code=400, detail="Email is required.")
    try:
        # Get or create user
        r = await db.execute(select(PathfinderUser).where(PathfinderUser.email == email))
        user = r.scalar_one_or_none()
        if not user:
            user = PathfinderUser(email=email, name=request.name or email.split("@")[0])
            db.add(user)
            await db.flush()
        else:
            if request.name:
                user.name = request.name

        # Upsert profile
        r = await db.execute(select(PathfinderUserProfile).where(PathfinderUserProfile.user_id == user.id))
        profile = r.scalar_one_or_none()
        if not profile:
            profile = PathfinderUserProfile(user_id=user.id)
            db.add(profile)
            await db.flush()
        profile.name = request.name or profile.name
        profile.academic_title = request.academic_title or ""
        profile.technical_skills = request.technical_skills or []
        profile.soft_skills = request.soft_skills or []

        # Replace coursework
        await db.execute(delete(PathfinderUserCoursework).where(PathfinderUserCoursework.user_id == user.id))
        for c in (request.courses or []):
            title = (c.get("title") or "").strip()
            if title:
                db.add(PathfinderUserCoursework(
                    user_id=user.id, title=title,
                    term=str(c.get("term", "—"))[:100],
                    grade=str(c.get("grade", "—"))[:50],
                    tags=c.get("tags") or [],
                ))

        # Replace projects
        await db.execute(delete(PathfinderUserProject).where(PathfinderUserProject.user_id == user.id))
        for p in (request.profile_projects or []):
            title = (p.get("title") or "").strip()
            if title:
                db.add(PathfinderUserProject(
                    user_id=user.id, title=title,
                    description=(p.get("description") or "")[:5000],
                    technologies=p.get("technologies") or [],
                    date=str(p.get("date", "—"))[:100],
                ))

        # Replace career interests
        await db.execute(delete(PathfinderUserCareerInterest).where(PathfinderUserCareerInterest.user_id == user.id))
        for i in (request.career_interests or []):
            interest = (i or "").strip()
            if interest:
                db.add(PathfinderUserCareerInterest(user_id=user.id, interest=interest))

        # Replace documents metadata (optional)
        await db.execute(delete(PathfinderUserDocument).where(PathfinderUserDocument.user_id == user.id))
        for d in (request.documents or []):
            cat = (d.get("category") or "other")[:50]
            fn = (d.get("filename") or "file")[:500]
            db.add(PathfinderUserDocument(
                user_id=user.id, category=cat, filename=fn,
                extracted_text=(d.get("extracted_text") or "")[:15000],
            ))

        await db.commit()
        return {"ok": True, "message": "Profile saved."}
    except HTTPException:
        raise
    except Exception as e:
        await db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/career/profile")
async def get_profile(
    email: str = Query(..., description="User email"),
    db: AsyncSession = Depends(get_db),
):
    """Load user profile from PostgreSQL."""
    email = (email or "").strip().lower()
    if not email:
        raise HTTPException(status_code=400, detail="Email is required.")
    try:
        r = await db.execute(select(PathfinderUser).where(PathfinderUser.email == email))
        user = r.scalar_one_or_none()
        if not user:
            return {"profile": None, "courses": [], "projects": [], "career_interests": []}

        r = await db.execute(select(PathfinderUserProfile).where(PathfinderUserProfile.user_id == user.id))
        profile = r.scalar_one_or_none()
        r = await db.execute(select(PathfinderUserCoursework).where(PathfinderUserCoursework.user_id == user.id))
        courses = [{"title": c.title, "term": c.term or "—", "grade": c.grade or "—", "tags": c.tags or []} for c in r.scalars().all()]
        r = await db.execute(select(PathfinderUserProject).where(PathfinderUserProject.user_id == user.id))
        projects = [{"title": p.title, "description": p.description or "", "technologies": p.technologies or [], "date": p.date or "—"} for p in r.scalars().all()]
        r = await db.execute(select(PathfinderUserCareerInterest).where(PathfinderUserCareerInterest.user_id == user.id))
        interests = [i.interest for i in r.scalars().all()]

        return {
            "profile": {
                "name": profile.name if profile else user.name,
                "academic_title": profile.academic_title if profile else "",
                "technical_skills": profile.technical_skills if profile else [],
                "soft_skills": profile.soft_skills if profile else [],
            },
            "courses": courses,
            "projects": projects,
            "career_interests": interests,
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


def _normalize_skill(s: str) -> str:
    return re.sub(r"[^a-z0-9]", "", (s or "").lower())


class SaveCareerInterestsRequest(BaseModel):
    email: str
    career_interests: List[str] = []


@router.put("/career/profile/career-interests")
async def save_career_interests(
    request: SaveCareerInterestsRequest,
    db: AsyncSession = Depends(get_db),
):
    """Update only career interests (call when user adds/removes interests)."""
    email = (request.email or "").strip().lower()
    interests = request.career_interests or []
    if not email:
        raise HTTPException(status_code=400, detail="Email is required.")
    try:
        r = await db.execute(select(PathfinderUser).where(PathfinderUser.email == email))
        user = r.scalar_one_or_none()
        if not user:
            user = PathfinderUser(email=email, name=email.split("@")[0])
            db.add(user)
            await db.flush()
        await db.execute(delete(PathfinderUserCareerInterest).where(PathfinderUserCareerInterest.user_id == user.id))
        for i in interests:
            interest = (i or "").strip()
            if interest:
                db.add(PathfinderUserCareerInterest(user_id=user.id, interest=interest))
        await db.commit()
        return {"ok": True}
    except HTTPException:
        raise
    except Exception as e:
        await db.rollback()
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/career/related-jobs")
async def get_related_jobs(
    email: str = Query(..., description="User email"),
    limit: int = Query(20, ge=1, le=50),
    db: AsyncSession = Depends(get_db),
):
    """Get jobs related to user profile (skills, projects, interests)."""
    email = (email or "").strip().lower()
    if not email:
        raise HTTPException(status_code=400, detail="Email is required.")
    try:
        # Load profile
        r = await db.execute(select(PathfinderUser).where(PathfinderUser.email == email))
        user = r.scalar_one_or_none()
        if not user:
            skills_set = set()
        else:
            r = await db.execute(select(PathfinderUserProfile).where(PathfinderUserProfile.user_id == user.id))
            profile = r.scalar_one_or_none()
            r = await db.execute(select(PathfinderUserProject).where(PathfinderUserProject.user_id == user.id))
            projects = r.scalars().all()
            r = await db.execute(select(PathfinderUserCareerInterest).where(PathfinderUserCareerInterest.user_id == user.id))
            interests = r.scalars().all()

            skills_set = set()
            for s in (profile.technical_skills or []) if profile else []:
                n = (s.get("name") if isinstance(s, dict) else str(s)) or ""
                if n:
                    skills_set.add(_normalize_skill(n))
            for p in projects:
                for t in (p.technologies or []):
                    if t:
                        skills_set.add(_normalize_skill(str(t)))
            for i in interests:
                if i.interest:
                    skills_set.add(_normalize_skill(i.interest))

        # Load all jobs and score by overlap
        r = await db.execute(select(PathfinderJob).limit(200))
        jobs = r.scalars().all()
        scored = []
        for j in jobs:
            job_skills = re.split(r"[,/;\s]+", (j.required_skills or "") + " " + (j.title or ""))
            job_norm = {_normalize_skill(x) for x in job_skills if x}
            overlap = len(skills_set & job_norm) if skills_set else 0
            scored.append((overlap, j))

        scored.sort(key=lambda x: (-x[0], x[1].title))
        result = []
        for overlap, j in scored[:limit]:
            # Match score: base 60 + up to 40 from skill overlap (5+ matches = 100%)
            match_score = min(100, 60 + overlap * 8)
            result.append({
                "id": j.id,
                "title": j.title,
                "company": j.company,
                "description": j.description,
                "required_skills": j.required_skills,
                "location": j.location,
                "job_type": j.job_type,
                "industry": j.industry,
                "salary": getattr(j, "salary", None),
                "match_score": match_score,
            })
        return {"jobs": result}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
